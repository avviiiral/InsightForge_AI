"""ExportAgent — renders a ReportAgent report (or raw DataFrame) to disk.

Supported formats: PDF, Excel (.xlsx), Markdown (.md), PowerPoint (.pptx),
and CSV (cleaned data only). Each `export_*` method returns the output
file path so callers (Streamlit download buttons, FastAPI file responses)
can stream the result back to the user.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
)

from app.agents.base_agent import BaseAgent


class ExportAgent(BaseAgent):
    name = "export_agent"
    description = "Exports an assembled report or dataframe to PDF, Excel, Markdown, PPTX, or CSV."

    def _execute(
        self,
        export_format: str,
        output_dir: str,
        filename_stem: str,
        report: dict[str, Any] | None = None,
        dataframe: pd.DataFrame | None = None,
    ) -> dict[str, Any]:
        out_dir = Path(output_dir)
        out_dir.mkdir(parents=True, exist_ok=True)

        dispatch = {
            "pdf": self.export_pdf,
            "excel": self.export_excel,
            "markdown": self.export_markdown,
            "pptx": self.export_pptx,
            "csv": self.export_csv,
        }
        if export_format not in dispatch:
            raise ValueError(f"Unsupported export_format '{export_format}'. Options: {list(dispatch)}")

        if export_format == "csv":
            if dataframe is None:
                raise ValueError("`dataframe` is required for CSV export.")
            path = self.export_csv(dataframe, out_dir, filename_stem)
        else:
            if report is None:
                raise ValueError(f"`report` is required for {export_format} export.")
            path = dispatch[export_format](report, out_dir, filename_stem)

        return {"file_path": str(path), "format": export_format}

    # -- CSV --------------------------------------------------------------

    @staticmethod
    def export_csv(dataframe: pd.DataFrame, out_dir: Path, stem: str) -> Path:
        path = out_dir / f"{stem}.csv"
        dataframe.to_csv(path, index=False)
        return path

    # -- Excel --------------------------------------------------------------

    @staticmethod
    def export_excel(report: dict[str, Any], out_dir: Path, stem: str) -> Path:
        path = out_dir / f"{stem}.xlsx"
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            summary_rows = [{"Metric": k, "Value": v} for k, v in report.get("data_quality", {}).items()
                             if not isinstance(v, (list, dict))]
            pd.DataFrame(summary_rows).to_excel(writer, sheet_name="Data Quality", index=False)

            pd.DataFrame(report.get("kpis", [])).to_excel(writer, sheet_name="KPIs", index=False)
            pd.DataFrame(report.get("insights", [])).to_excel(writer, sheet_name="Insights", index=False)
            pd.DataFrame(report.get("recommendations", [])).to_excel(writer, sheet_name="Recommendations", index=False)

            corr_matrix = report.get("correlation", {}).get("matrix")
            if corr_matrix:
                pd.DataFrame(corr_matrix).to_excel(writer, sheet_name="Correlation Matrix")

            profiles = report.get("column_profiles", {})
            if profiles:
                rows = [{"column": col, **info} for col, info in profiles.items()]
                pd.json_normalize(rows).to_excel(writer, sheet_name="Column Profiles", index=False)
        return path

    # -- Markdown --------------------------------------------------------------

    @staticmethod
    def export_markdown(report: dict[str, Any], out_dir: Path, stem: str) -> Path:
        path = out_dir / f"{stem}.md"
        lines = [f"# {report.get('title', 'Analytics Report')}", "", f"_Generated: {report.get('generated_at')}_", ""]

        lines.append("## Executive Summary")
        for bullet in report.get("executive_summary", {}).get("bullets", []):
            lines.append(f"- {bullet}")
        lines.append("")

        if report.get("business_narrative"):
            lines += ["## Business Narrative", report["business_narrative"], ""]

        quality = report.get("data_quality", {})
        lines += [
            "## Data Quality",
            f"- Health score: **{quality.get('health_score')}/100**",
            f"- Rows: {quality.get('total_rows')}, Columns: {quality.get('total_columns')}",
            f"- Duplicate rows: {quality.get('duplicate_rows')}",
            f"- Missing cells: {quality.get('total_missing_cells')} ({quality.get('missing_pct')}%)",
            "",
        ]

        lines.append("## KPIs")
        for kpi in report.get("kpis", []):
            lines.append(f"- **{kpi['name']}**: {kpi['formatted_value']}")
        lines.append("")

        lines.append("## Key Insights")
        for insight in report.get("insights", []):
            lines.append(f"- **[{insight.get('severity', 'info').upper()}] {insight.get('title')}**: {insight.get('description')}")
        lines.append("")

        lines.append("## Recommendations")
        for rec in report.get("recommendations", []):
            lines.append(f"- **{rec.get('title')}**: {rec.get('action')}")
        lines.append("")

        path.write_text("\n".join(lines), encoding="utf-8")
        return path

    # -- PDF --------------------------------------------------------------

    @staticmethod
    def export_pdf(report: dict[str, Any], out_dir: Path, stem: str) -> Path:
        path = out_dir / f"{stem}.pdf"
        doc = SimpleDocTemplate(str(path), pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle("TitleCustom", parent=styles["Title"], textColor=colors.HexColor("#1e3a8a"))
        heading_style = ParagraphStyle("HeadingCustom", parent=styles["Heading2"], textColor=colors.HexColor("#1e3a8a"),
                                       spaceBefore=14, spaceAfter=6)
        body_style = styles["BodyText"]

        story: list[Any] = [Paragraph(report.get("title", "Analytics Report"), title_style),
                             Paragraph(f"Generated: {report.get('generated_at')}", styles["Normal"]), Spacer(1, 16)]

        story.append(Paragraph("Executive Summary", heading_style))
        for bullet in report.get("executive_summary", {}).get("bullets", []):
            story.append(Paragraph(f"• {bullet}", body_style))

        if report.get("business_narrative"):
            story.append(Paragraph("Business Narrative", heading_style))
            story.append(Paragraph(report["business_narrative"], body_style))

        quality = report.get("data_quality", {})
        story.append(Paragraph("Data Quality", heading_style))
        quality_table_data = [
            ["Metric", "Value"],
            ["Health Score", f"{quality.get('health_score')}/100"],
            ["Total Rows", str(quality.get("total_rows"))],
            ["Total Columns", str(quality.get("total_columns"))],
            ["Duplicate Rows", str(quality.get("duplicate_rows"))],
            ["Missing Cells", f"{quality.get('total_missing_cells')} ({quality.get('missing_pct')}%)"],
        ]
        story.append(ExportAgent._styled_table(quality_table_data))

        story.append(Paragraph("KPIs", heading_style))
        kpi_data = [["KPI", "Value"]] + [[k["name"], k["formatted_value"]] for k in report.get("kpis", [])]
        if len(kpi_data) > 1:
            story.append(ExportAgent._styled_table(kpi_data))

        story.append(PageBreak())
        story.append(Paragraph("Key Insights", heading_style))
        for insight in report.get("insights", []):
            story.append(Paragraph(f"<b>[{insight.get('severity', 'info').upper()}] {insight.get('title')}</b>", body_style))
            story.append(Paragraph(insight.get("description", ""), body_style))
            story.append(Spacer(1, 4))

        story.append(Paragraph("Recommendations", heading_style))
        for rec in report.get("recommendations", []):
            story.append(Paragraph(f"<b>{rec.get('title')}</b>: {rec.get('action')}", body_style))
            story.append(Spacer(1, 4))

        doc.build(story)
        return path

    @staticmethod
    def _styled_table(data: list[list[str]]) -> Table:
        table = Table(data, hAlign="LEFT", colWidths=[200, 250])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e3a8a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1f5f9")]),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return table

    # -- PowerPoint --------------------------------------------------------------

    @staticmethod
    def export_pptx(report: dict[str, Any], out_dir: Path, stem: str) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        path = out_dir / f"{stem}.pptx"
        prs = Presentation()
        navy = RGBColor(0x1E, 0x3A, 0x8A)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = report.get("title", "Analytics Report")
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = navy
        slide.placeholders[1].text = f"Generated: {report.get('generated_at')}"

        # Executive summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Executive Summary"
        body = slide.placeholders[1].text_frame
        bullets = report.get("executive_summary", {}).get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)

        # KPI slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Key Performance Indicators"
        body = slide.placeholders[1].text_frame
        kpis = report.get("kpis", [])
        for i, kpi in enumerate(kpis):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = f"{kpi['name']}: {kpi['formatted_value']}"
            p.font.size = Pt(18)

        # Insights slide(s)
        insights = report.get("insights", [])
        for chunk_start in range(0, len(insights), 5):
            chunk = insights[chunk_start: chunk_start + 5]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = "Key Insights"
            body = slide.placeholders[1].text_frame
            for i, insight in enumerate(chunk):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = f"[{insight.get('severity', 'info').upper()}] {insight.get('title')}"
                p.font.size = Pt(16)

        # Recommendations slide
        recs = report.get("recommendations", [])
        if recs:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = "Recommendations"
            body = slide.placeholders[1].text_frame
            for i, rec in enumerate(recs[:6]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = f"{rec.get('title')}: {rec.get('action')}"
                p.font.size = Pt(14)

        prs.save(str(path))
        return path
